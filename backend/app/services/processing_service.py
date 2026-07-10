"""
Processing service — job queue orchestrator + metadata extraction pipeline.


Architecture:
  - enqueue_job()   → inserts a ResourceProcessingJob row
  - poll_job_queue() → async loop, runs forever, polls every 5s
  - _run_job()       → dispatches to correct extractor
  - _extract_metadata() → downloads file, extracts metadata, writes resource_metadata

Future extensions: add job_type == "chunking", "embedding" etc — each gets its own
dispatcher branch. The worker loop does NOT need to change.
"""

import asyncio
import io
import mimetypes
from datetime import datetime, timezone
from uuid import UUID, uuid4

import structlog
from sqlalchemy import select, update
from sqlalchemy.ext.asyncio import AsyncSession

from app.db.session import AsyncSessionLocal
from app.models.resource import Resource, ResourceChunk, ResourceMetadata, ResourceProcessingJob
from app.repositories.chunk_repository import ChunkRepository
from app.services import chunking_service, embedding_service, extraction_service, storage_service

logger = structlog.get_logger()

_worker_running = False


# ─────────────────────────────────────────────────────────────────────────────
# Queue management
# ─────────────────────────────────────────────────────────────────────────────

async def enqueue_job(
    db: AsyncSession,
    resource_id: UUID,
    job_type: str = "metadata_extraction",
) -> ResourceProcessingJob:
    """Insert a new job row. Called inside an existing transaction."""
    job = ResourceProcessingJob(
        resource_id=resource_id,
        job_type=job_type,
        status="queued",
        attempts=0,
        max_attempts=3,
    )
    db.add(job)
    await db.flush()
    logger.info("processing.job_enqueued", resource_id=str(resource_id), job_type=job_type)
    return job


async def poll_job_queue() -> None:
    """
    Async background loop. Started once at app startup via lifespan.
    Polls the resource_processing_jobs table every 5 seconds.
    """
    global _worker_running
    _worker_running = True
    logger.info("processing.worker_started")

    while True:
        try:
            await _process_next_batch()
        except Exception as e:
            logger.error("processing.poll_error", error=str(e))
        await asyncio.sleep(5)


async def _process_next_batch() -> None:
    """Pick up to 3 queued jobs and process them concurrently."""
    async with AsyncSessionLocal() as db:
        result = await db.execute(
            select(ResourceProcessingJob)
            .where(ResourceProcessingJob.status == "queued")
            .order_by(ResourceProcessingJob.created_at.asc())
            .limit(3)
            .with_for_update(skip_locked=True)
        )
        jobs = result.scalars().all()
        if not jobs:
            return

        # Mark all as running before releasing lock
        for job in jobs:
            job.status = "running"
            job.started_at = datetime.now(timezone.utc)
            job.attempts += 1
            job.worker_id = f"worker-{uuid4().hex[:8]}"
        await db.commit()

    # Process outside the lock
    await asyncio.gather(*[_run_job(job.id) for job in jobs], return_exceptions=True)


async def _run_job(job_id: UUID) -> None:
    """Dispatch a single job to the correct handler."""
    async with AsyncSessionLocal() as db:
        job_res = await db.execute(
            select(ResourceProcessingJob).where(ResourceProcessingJob.id == job_id)
        )
        job = job_res.scalar_one_or_none()
        if not job:
            return

        resource_res = await db.execute(
            select(Resource).where(Resource.id == job.resource_id)
        )
        resource = resource_res.scalar_one_or_none()
        if not resource:
            await _fail_job(db, job, "Resource not found")
            return

        try:
            if job.job_type == "metadata_extraction":
                await _extract_metadata(db, resource, job)
            elif job.job_type == "chunking":
                await _run_chunking(db, resource, job)
            elif job.job_type == "embedding":
                await _run_embedding(db, resource, job)
            else:
                logger.warning("processing.unknown_job_type", job_type=job.job_type)
                await _fail_job(db, job, f"Unknown job_type: {job.job_type}")
        except Exception as e:
            logger.error("processing.job_failed", job_id=str(job_id), error=str(e))
            if job.attempts >= job.max_attempts:
                await _fail_job(db, job, str(e))
            else:
                # Return to queue for retry
                job.status = "queued"
                job.error_message = str(e)
                await db.commit()


async def _fail_job(db: AsyncSession, job: ResourceProcessingJob, error: str) -> None:
    job.status = "failed"
    job.error_message = error
    job.completed_at = datetime.now(timezone.utc)

    # Update resource stage to failed
    await db.execute(
        update(Resource)
        .where(Resource.id == job.resource_id)
        .values(
            processing_stage="failed",
            processing_status="failed",
            processing_error=error,
        )
    )
    await db.commit()
    logger.error("processing.job_permanent_failure", job_id=str(job.id), error=error)


# ─────────────────────────────────────────────────────────────────────────────
# Metadata extraction pipeline
# ─────────────────────────────────────────────────────────────────────────────

async def _set_stage(db: AsyncSession, resource: Resource, stage: str) -> None:
    resource.processing_stage = stage
    await db.execute(
        update(Resource)
        .where(Resource.id == resource.id)
        .values(processing_stage=stage)
    )
    await db.commit()
    logger.info("processing.stage_changed", resource_id=str(resource.id), stage=stage)


async def _extract_metadata(
    db: AsyncSession,
    resource: Resource,
    job: ResourceProcessingJob,
) -> None:
    rid = resource.id
    path = resource.file_url

    # ── Stage 1: Validate file exists ────────────────────────────────────
    await _set_stage(db, resource, "validating")
    exists = await storage_service.verify_object_exists(path)
    if not exists:
        raise RuntimeError(f"File not found in storage at path: {path}")

    # ── Stage 2: Download + extract ───────────────────────────────────────
    await _set_stage(db, resource, "extracting")
    raw_bytes = await storage_service.download_object(path)
    file_type = resource.file_type.lower()
    meta: dict = {}

    if file_type == "pdf":
        meta = await _extract_pdf(raw_bytes)
    elif file_type in ("docx", "doc"):
        meta = await _extract_docx(raw_bytes)
    elif file_type in ("pptx", "ppt"):
        meta = await _extract_pptx(raw_bytes)
    elif file_type in ("txt", "md"):
        meta = _extract_text(raw_bytes)
    elif file_type in ("jpg", "jpeg", "png", "gif", "webp", "svg"):
        meta = await _extract_image(raw_bytes, file_type)
    else:
        meta = {"raw": {"file_type": file_type, "note": "no extractor"}}

    # ── Write resource_metadata row ───────────────────────────────────────
    existing = await db.execute(
        select(ResourceMetadata).where(ResourceMetadata.resource_id == rid)
    )
    rm = existing.scalar_one_or_none()
    if rm:
        for k, v in meta.items():
            if k != "raw":
                setattr(rm, k, v)
        rm.raw = meta.get("raw", {})
        rm.extracted_at = datetime.now(timezone.utc)
    else:
        raw_extra = meta.pop("raw", {})
        rm = ResourceMetadata(resource_id=rid, raw=raw_extra, **meta)
        db.add(rm)

    await db.flush()

    # ── Metadata done → hand off to chunking ──────────────────────────────
    now = datetime.now(timezone.utc)
    await db.execute(
        update(Resource).where(Resource.id == rid).values(processing_status="processing")
    )
    job.status = "completed"
    job.completed_at = now
    await enqueue_job(db, rid, "chunking")
    await db.commit()
    logger.info("processing.metadata_done", resource_id=str(rid))


# ─────────────────────────────────────────────────────────────────────────────
# Chunking pipeline
# ─────────────────────────────────────────────────────────────────────────────

async def _run_chunking(db: AsyncSession, resource: Resource, job: ResourceProcessingJob) -> None:
    """Extract full text, split into chunks, persist them (no embeddings yet)."""
    await _set_stage(db, resource, "chunking")
    raw_bytes = await storage_service.download_object(resource.file_url)

    try:
        blocks = extraction_service.extract_text_blocks(raw_bytes, resource.file_type)
    except extraction_service.UnsupportedFileType:
        # e.g. images — nothing to chunk. Complete, but not AI-ready.
        await _complete_without_ai(db, resource, job, reason="no_text_extractor")
        return

    chunk_data = chunking_service.chunk_blocks(blocks)
    if not chunk_data:
        await _complete_without_ai(db, resource, job, reason="no_extractable_text")
        return

    repo = ChunkRepository(db)
    await repo.delete_chunks(resource.id)  # idempotent — supports stage re-run
    rows = [
        ResourceChunk(
            resource_id=resource.id,
            vault_id=resource.vault_id,
            chunk_index=c.chunk_index,
            content=c.content,
            token_count=c.token_count,
            page_number=c.page_number,
            heading=c.heading,
        )
        for c in chunk_data
    ]
    await repo.bulk_insert(rows)

    job.status = "completed"
    job.completed_at = datetime.now(timezone.utc)
    await enqueue_job(db, resource.id, "embedding")
    await db.commit()
    logger.info("processing.chunking_done", resource_id=str(resource.id), chunks=len(rows))


async def _run_embedding(db: AsyncSession, resource: Resource, job: ResourceProcessingJob) -> None:
    """Embed all chunks; only mark AI-ready once every chunk has a vector."""
    await _set_stage(db, resource, "embedding")
    await embedding_service.embed_resource_chunks(db, resource)

    repo = ChunkRepository(db)
    total = await repo.count_for_resource(resource.id)
    embedded = await repo.count_embedded(resource.id)
    if total == 0 or embedded < total:
        # Incomplete → raise so the worker retries the embedding stage only.
        raise RuntimeError(f"Embedding incomplete: {embedded}/{total} chunks embedded.")

    now = datetime.now(timezone.utc)
    await db.execute(
        update(Resource)
        .where(Resource.id == resource.id)
        .values(
            processing_stage="complete",
            processing_status="ready",
            is_ai_ready=True,
            processing_error=None,
            processed_at=now,
        )
    )
    job.status = "completed"
    job.completed_at = now
    await db.commit()
    logger.info("processing.embedding_done", resource_id=str(resource.id), embedded=embedded)


async def _complete_without_ai(
    db: AsyncSession, resource: Resource, job: ResourceProcessingJob, *, reason: str
) -> None:
    now = datetime.now(timezone.utc)
    await db.execute(
        update(Resource)
        .where(Resource.id == resource.id)
        .values(
            processing_stage="complete",
            processing_status="ready",
            is_ai_ready=False,
            processing_error=None,
            processed_at=now,
        )
    )
    job.status = "completed"
    job.completed_at = now
    await db.commit()
    logger.info("processing.complete_without_ai", resource_id=str(resource.id), reason=reason)


# ─────────────────────────────────────────────────────────────────────────────
# Extractors
# ─────────────────────────────────────────────────────────────────────────────

async def _extract_pdf(data: bytes) -> dict:
    try:
        import pypdf
        reader = pypdf.PdfReader(io.BytesIO(data))
        info = reader.metadata or {}
        pages = len(reader.pages)

        # Estimate word count from first 10 pages (avoid full scan for large PDFs)
        words = 0
        for page in reader.pages[:10]:
            text = page.extract_text() or ""
            words += len(text.split())
        if pages > 10:
            words = int(words / 10 * pages)  # extrapolate

        return {
            "pages": pages,
            "words": words,
            "author": str(info.get("/Author", "")) or None,
            "detected_title": str(info.get("/Title", "")) or None,
            "pdf_version": reader.pdf_header.replace("%PDF-", "") if reader.pdf_header else None,
            "reading_time_mins": max(1, words // 200),
            "raw": {"creator": str(info.get("/Creator", "")), "producer": str(info.get("/Producer", ""))},
        }
    except Exception as e:
        logger.warning("extractor.pdf_failed", error=str(e))
        return {"raw": {"error": str(e)}}


async def _extract_docx(data: bytes) -> dict:
    try:
        import docx
        doc = docx.Document(io.BytesIO(data))
        words = sum(len(p.text.split()) for p in doc.paragraphs)
        core = doc.core_properties
        return {
            "words": words,
            "author": core.author or None,
            "detected_title": core.title or None,
            "reading_time_mins": max(1, words // 200),
            "raw": {"revision": core.revision},
        }
    except Exception as e:
        logger.warning("extractor.docx_failed", error=str(e))
        return {"raw": {"error": str(e)}}


async def _extract_pptx(data: bytes) -> dict:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        slides = len(prs.slides)
        words = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    words += len(shape.text.split())
        core = prs.core_properties
        return {
            "slides": slides,
            "words": words,
            "author": core.author or None,
            "detected_title": core.title or None,
            "reading_time_mins": max(1, slides * 2),
            "raw": {},
        }
    except Exception as e:
        logger.warning("extractor.pptx_failed", error=str(e))
        return {"raw": {"error": str(e)}}


def _extract_text(data: bytes) -> dict:
    try:
        text = data.decode("utf-8", errors="replace")
        words = len(text.split())
        lines = text.count("\n")
        return {
            "words": words,
            "reading_time_mins": max(1, words // 200),
            "raw": {"lines": lines, "chars": len(text)},
        }
    except Exception as e:
        return {"raw": {"error": str(e)}}


async def _extract_image(data: bytes, file_type: str) -> dict:
    try:
        from PIL import Image
        img = Image.open(io.BytesIO(data))
        return {
            "width_px": img.width,
            "height_px": img.height,
            "raw": {"format": img.format, "mode": img.mode},
        }
    except Exception as e:
        logger.warning("extractor.image_failed", error=str(e))
        return {"raw": {"error": str(e)}}
