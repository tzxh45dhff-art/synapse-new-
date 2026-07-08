"""ExtractionService — turn raw file bytes into structured text blocks.

Returns a list of :class:`ExtractedBlock`, each carrying the text plus its
page number and heading where the format exposes them. The chunking service
consumes these; page/heading provenance flows all the way to citations.

Extensible by design: register a new ``(file_type) -> extractor`` entry and the
rest of the pipeline is unchanged. Images/OCR/audio/video plug in the same way.
"""

from __future__ import annotations

import io
from dataclasses import dataclass
from typing import Callable

import structlog

logger = structlog.get_logger()


@dataclass(slots=True)
class ExtractedBlock:
    text: str
    page_number: int | None = None
    heading: str | None = None


class UnsupportedFileType(Exception):
    pass


# ─────────────────────────────────────────────────────────────────────────────
# Public API
# ─────────────────────────────────────────────────────────────────────────────

def extract_text_blocks(raw: bytes, file_type: str) -> list[ExtractedBlock]:
    """Extract structured text blocks from a document.

    Raises :class:`UnsupportedFileType` for formats with no text extractor
    (e.g. images — reserved for a future OCR extractor).
    """
    ft = file_type.lower()
    extractor = _EXTRACTORS.get(ft)
    if not extractor:
        raise UnsupportedFileType(f"No text extractor for file_type '{file_type}'.")
    blocks = extractor(raw)
    # Drop empty/whitespace-only blocks.
    return [b for b in blocks if b.text and b.text.strip()]


# ─────────────────────────────────────────────────────────────────────────────
# Extractors
# ─────────────────────────────────────────────────────────────────────────────

def _extract_pdf(raw: bytes) -> list[ExtractedBlock]:
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(raw))
    blocks: list[ExtractedBlock] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            blocks.append(ExtractedBlock(text=text, page_number=i))
    return blocks


def _extract_docx(raw: bytes) -> list[ExtractedBlock]:
    import docx

    doc = docx.Document(io.BytesIO(raw))
    blocks: list[ExtractedBlock] = []
    current_heading: str | None = None
    buffer: list[str] = []

    def flush() -> None:
        if buffer:
            blocks.append(ExtractedBlock(text="\n".join(buffer), heading=current_heading))
            buffer.clear()

    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        if style.startswith("heading") and para.text.strip():
            flush()
            current_heading = para.text.strip()
            buffer.append(para.text)
        elif para.text.strip():
            buffer.append(para.text)
    flush()

    # Tables → markdown-ish rows so the chunker keeps them intact.
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            blocks.append(ExtractedBlock(text="\n".join(rows), heading=current_heading))

    return blocks


def _extract_pptx(raw: bytes) -> list[ExtractedBlock]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    blocks: list[ExtractedBlock] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = None
        parts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            txt = shape.text.strip()
            if not txt:
                continue
            if title is None and shape == getattr(slide.shapes, "title", None):
                title = txt
            parts.append(txt)
        if parts:
            blocks.append(ExtractedBlock(text="\n".join(parts), page_number=i, heading=title))
    return blocks


def _extract_markdown(raw: bytes) -> list[ExtractedBlock]:
    text = raw.decode("utf-8", errors="replace")
    blocks: list[ExtractedBlock] = []
    current_heading: str | None = None
    buffer: list[str] = []
    in_code = False

    def flush() -> None:
        if buffer:
            blocks.append(ExtractedBlock(text="\n".join(buffer), heading=current_heading))
            buffer.clear()

    for line in text.splitlines():
        stripped = line.strip()
        if stripped.startswith("```"):
            in_code = not in_code
            buffer.append(line)
            continue
        if not in_code and stripped.startswith("#"):
            flush()
            current_heading = stripped.lstrip("#").strip()
            buffer.append(line)
        else:
            buffer.append(line)
    flush()
    return blocks


def _extract_txt(raw: bytes) -> list[ExtractedBlock]:
    text = raw.decode("utf-8", errors="replace")
    return [ExtractedBlock(text=text)]


_EXTRACTORS: dict[str, Callable[[bytes], list[ExtractedBlock]]] = {
    "pdf": _extract_pdf,
    "docx": _extract_docx,
    "doc": _extract_docx,
    "pptx": _extract_pptx,
    "ppt": _extract_pptx,
    "md": _extract_markdown,
    "markdown": _extract_markdown,
    "txt": _extract_txt,
}
